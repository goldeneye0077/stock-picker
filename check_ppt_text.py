"""
检查生成的 PPT 中提取的文字
"""
from pptx import Presentation

prs = Presentation("演示文稿1_提取文字.pptx")

print("检查前 3 张幻灯片的文字内容：\n")

for i in range(min(3, len(prs.slides))):
    slide = prs.slides[i]
    print(f"=== 第 {i+1} 张幻灯片 ===")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            print(f"文字内容: {shape.text}")
            print()
