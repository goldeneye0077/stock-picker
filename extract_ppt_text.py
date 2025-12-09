"""
从PPT图片中提取文字并生成新的PPT
使用 PaddleOCR 进行中文OCR识别
"""
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from paddleocr import PaddleOCR

# 初始化 PaddleOCR（使用新版API）
ocr = PaddleOCR(use_textline_orientation=True, lang='ch')

def extract_images_from_ppt(ppt_path):
    """从PPT中提取所有图片"""
    prs = Presentation(ppt_path)
    images = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_images = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # 获取图片二进制数据
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                slide_images.append(image)
                print(f"幻灯片 {slide_idx + 1}: 找到 1 张图片")

        images.append(slide_images)

    return images

def ocr_image(image):
    """使用OCR识别图片中的文字"""
    # 将PIL Image转换为临时文件路径
    temp_path = "temp_image.png"
    image.save(temp_path)

    try:
        # 使用PaddleOCR的新版 predict 方法识别
        result = ocr.predict(temp_path)

        # 提取文字
        texts = []
        if result and isinstance(result, list) and len(result) > 0:
            # result 是一个列表，第一个元素是字典
            first_result = result[0]
            if isinstance(first_result, dict):
                # 新版API返回格式：rec_texts 是文字列表，rec_scores 是置信度列表
                rec_texts = first_result.get('rec_texts', [])
                rec_scores = first_result.get('rec_scores', [])

                for i, text in enumerate(rec_texts):
                    confidence = rec_scores[i] if i < len(rec_scores) else 1.0
                    if confidence > 0.5:  # 只保留置信度>0.5的结果
                        texts.append(str(text))

        return '\n'.join(texts)
    finally:
        # 清理临时文件
        if os.path.exists(temp_path):
            os.remove(temp_path)

def create_new_ppt(original_ppt_path, output_path):
    """创建包含提取文字的新PPT（保留原始版式和内容）"""
    # 读取原始PPT
    original_prs = Presentation(original_ppt_path)
    images_by_slide = extract_images_from_ppt(original_ppt_path)

    # 复制原始PPT的所有幻灯片
    new_prs = Presentation(original_ppt_path)

    print("\n开始OCR识别...")

    # 处理每一页
    for slide_idx, slide_images in enumerate(images_by_slide):
        print(f"\n处理幻灯片 {slide_idx + 1}/{len(images_by_slide)}...")

        # 获取对应的幻灯片（已经是原始PPT的副本）
        slide = new_prs.slides[slide_idx]

        # 处理该页的所有图片
        all_texts = []
        for img_idx, image in enumerate(slide_images):
            print(f"  识别图片 {img_idx + 1}/{len(slide_images)}...")
            text = ocr_image(image)
            if text:
                all_texts.append(text)

        # 将提取的文字添加到幻灯片
        if all_texts:
            # 添加文本框
            left = Inches(0.5)
            top = Inches(0.5)
            width = new_prs.slide_width - Inches(1)
            height = new_prs.slide_height - Inches(1)

            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            # 添加文字
            combined_text = '\n\n'.join(all_texts)
            p = text_frame.paragraphs[0]
            p.text = combined_text
            p.font.size = Pt(14)
            p.font.name = '微软雅黑'

            print(f"  [成功] 提取到 {len(all_texts)} 段文字")
        else:
            # 如果没有识别到文字，添加提示
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(1)

            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = "（未识别到文字）"
            p.font.size = Pt(14)
            p.font.name = '微软雅黑'

            print(f"  [跳过] 未识别到文字")

    # 保存新PPT
    new_prs.save(output_path)
    print(f"\n[完成] 新PPT已保存到: {output_path}")

if __name__ == "__main__":
    input_ppt = "演示文稿1.pptx"
    output_ppt = "演示文稿1_提取文字.pptx"

    print("=" * 60)
    print("PPT图片文字提取工具")
    print("=" * 60)
    print(f"\n输入文件: {input_ppt}")
    print(f"输出文件: {output_ppt}")
    print("\n" + "=" * 60)

    try:
        create_new_ppt(input_ppt, output_ppt)
        print("\n" + "=" * 60)
        print("[成功] 处理完成！")
        print("=" * 60)
    except Exception as e:
        print(f"\n[错误] {e}")
        import traceback
        traceback.print_exc()
